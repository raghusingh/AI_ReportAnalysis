"""
AI Dashboard Studio — Config-Driven Edition
Reads dashboard_prompt.yaml + rules.yaml to determine what to render.
"""

import streamlit as st
import pandas as pd
import json, os, re, yaml, copy
from pathlib import Path
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import plotly.express as px
import plotly.graph_objects as go
from io import BytesIO
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import numpy as np

# ═══════════════════════════════════════════════════════════════
#  CONFIG LOADING
# ═══════════════════════════════════════════════════════════════
BASE_DIR = Path(__file__).parent

def load_yaml(filename):
    path = BASE_DIR / filename
    if path.exists():
        with open(path, "r", encoding="utf-8") as f:
            return yaml.safe_load(f) or {}
    return {}

prompt_cfg = load_yaml("dashboard_prompt.yaml")
rules_cfg  = load_yaml("rules.yaml")
dash_meta  = prompt_cfg.get("dashboard", {})

# ═══════════════════════════════════════════════════════════════
#  PAGE CONFIG
# ═══════════════════════════════════════════════════════════════
st.set_page_config(
    page_title=dash_meta.get("title", "AI Dashboard Studio"),
    page_icon=dash_meta.get("logo_emoji", "📊"),
    layout="wide",
    initial_sidebar_state="expanded",
)

# ═══════════════════════════════════════════════════════════════
#  CSS
# ═══════════════════════════════════════════════════════════════
st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Syne:wght@400;600;700;800&family=DM+Sans:wght@300;400;500&display=swap');
html,body,[class*="css"]{font-family:'DM Sans',sans-serif;}
.stApp{background:#f5f6fa;color:#1a1a2e;}
section[data-testid="stSidebar"]{background:#ffffff!important;border-right:1px solid #e2e4ed;box-shadow:2px 0 12px rgba(0,0,0,0.06);}
.main .block-container{padding:1.2rem 2rem;max-width:100%;}
.metric-tile{background:#fff;border:1px solid #e2e4ed;border-radius:12px;padding:1rem 1.2rem;position:relative;overflow:hidden;margin-bottom:.8rem;box-shadow:0 2px 8px rgba(109,40,217,.07);transition:box-shadow .2s;}
.metric-tile:hover{box-shadow:0 4px 16px rgba(109,40,217,.13);}
.metric-tile::before{content:'';position:absolute;top:0;left:0;width:4px;height:100%;background:linear-gradient(180deg,#7c3aed,#a78bfa);}
.metric-tile .label{font-size:.7rem;font-weight:600;text-transform:uppercase;letter-spacing:.1em;color:#8b8ba8;margin-bottom:.3rem;}
.metric-tile .value{font-family:'Syne',sans-serif;font-size:1.8rem;font-weight:800;color:#1a1a2e;line-height:1;}
.metric-tile .delta{font-size:.72rem;margin-top:.35rem;color:#059669;font-weight:500;}
.metric-tile .na{font-size:.8rem;color:#c4b5fd;font-style:italic;margin-top:.3rem;}
.section-title{font-family:'Syne',sans-serif;font-size:.85rem;font-weight:700;color:#6d28d9;text-transform:uppercase;letter-spacing:.09em;margin:1.2rem 0 .6rem;padding-bottom:.3rem;border-bottom:2px solid #ede9fe;}
.status-card{background:#fff;border:1px solid #e2e4ed;border-radius:12px;padding:1rem 1.2rem;margin-bottom:.8rem;box-shadow:0 2px 8px rgba(0,0,0,.05);}
.status-card .proj-name{font-family:'Syne',sans-serif;font-weight:700;font-size:.95rem;color:#1a1a2e;margin-bottom:.4rem;}
.status-badge{display:inline-block;padding:.2rem .7rem;border-radius:20px;font-size:.72rem;font-weight:700;text-transform:uppercase;letter-spacing:.06em;}
.badge-green{background:#dcfce7;color:#166534;}
.badge-amber{background:#fef9c3;color:#854d0e;}
.badge-red{background:#fee2e2;color:#991b1b;}
.badge-grey{background:#f3f4f6;color:#374151;}
.insight-card{background:#fff;border:1px solid #e2e4ed;border-radius:10px;padding:.9rem 1.1rem;margin-bottom:.6rem;font-size:.86rem;color:#4a4a6a;line-height:1.65;box-shadow:0 1px 4px rgba(0,0,0,.04);}
.insight-card .ititle{font-family:'Syne',sans-serif;font-weight:700;color:#1a1a2e;margin-bottom:.25rem;font-size:.9rem;}
.chat-user{background:#f0ebff;border-left:3px solid #7c3aed;color:#3b1f6e;padding:.8rem 1rem;border-radius:8px;margin-bottom:.5rem;font-size:.88rem;}
.chat-ai{background:#ecfdf5;border-left:3px solid #059669;color:#064e3b;padding:.8rem 1rem;border-radius:8px;margin-bottom:.5rem;font-size:.88rem;}
.stButton>button{background:linear-gradient(135deg,#6d28d9,#7c3aed)!important;color:white!important;border:none!important;border-radius:8px!important;font-family:'Syne',sans-serif!important;font-weight:700!important;padding:.45rem 1rem!important;box-shadow:0 2px 8px rgba(109,40,217,.25)!important;transition:opacity .2s!important;}
.stButton>button:hover{opacity:.88!important;}
.stTabs [data-baseweb="tab-list"]{background:#ede9fe;border-radius:10px;gap:4px;padding:4px;}
.stTabs [data-baseweb="tab"]{background:transparent;color:#7b78a0;border-radius:7px;font-family:'Syne',sans-serif;font-size:.82rem;font-weight:600;}
.stTabs [aria-selected="true"]{background:#fff!important;color:#6d28d9!important;box-shadow:0 1px 6px rgba(109,40,217,.15)!important;}
.stChatInput>div{background:#fff!important;border:1px solid #d4cef5!important;border-radius:10px!important;}
hr{border-color:#e2e4ed;}
::-webkit-scrollbar{width:5px;height:5px;}
::-webkit-scrollbar-thumb{background:#c4b5fd;border-radius:3px;}
.placeholder-box{background:#faf8ff;border:1px dashed #c4b5fd;border-radius:8px;padding:1.2rem;text-align:center;color:#a78bfa;font-size:.82rem;font-style:italic;margin-bottom:.6rem;}
</style>
""", unsafe_allow_html=True)

# ═══════════════════════════════════════════════════════════════
#  SESSION STATE
# ═══════════════════════════════════════════════════════════════
def init_session():
    defaults = {
        "api_key": "", "chat_history": [],
        "page_data": {}, "dataframes": {}, "sheet_names": [],
        "data_summary": "", "file_uploaded": False, "dashboard_ready": False,
        "live_pages": {},
    }
    for k, v in defaults.items():
        if k not in st.session_state:
            st.session_state[k] = v

init_session()

# ═══════════════════════════════════════════════════════════════
#  GEMINI
# ═══════════════════════════════════════════════════════════════
def get_model():
    if not st.session_state.api_key:
        return None
    model_name = rules_cfg.get("gemini", {}).get("model", "gemini-2.5-flash-lite-preview-06-17")
    genai.configure(api_key=st.session_state.api_key)
    return genai.GenerativeModel(model_name)

def call_gemini(prompt, system=""):
    model = get_model()
    if not model:
        return ""
    full = f"{system}\n\n{prompt}" if system else prompt
    try:
        response = model.generate_content(full)
        if hasattr(response, 'text') and response.text:
            return response.text
        if hasattr(response, 'candidates') and response.candidates:
            parts = response.candidates[0].content.parts
            return "".join(p.text for p in parts if hasattr(p, 'text'))
        return ""
    except Exception as e:
        return f"ERROR: {e}"

def parse_json(raw):
    if not raw:
        return {}
    raw = re.sub(r'<thinking>.*?</thinking>', '', raw, flags=re.DOTALL)
    raw = raw.strip()
    raw = re.sub(r'^```json\s*', '', raw, flags=re.MULTILINE)
    raw = re.sub(r'^```\s*', '', raw, flags=re.MULTILINE)
    raw = re.sub(r'```$', '', raw, flags=re.MULTILINE)
    raw = raw.strip()
    try:
        return json.loads(raw)
    except Exception:
        pass
    m = re.search(r'\{.*\}', raw, re.DOTALL)
    if m:
        try:
            return json.loads(m.group())
        except Exception:
            pass
    m = re.search(r'\[.*\]', raw, re.DOTALL)
    if m:
        try:
            return json.loads(m.group())
        except Exception:
            pass
    return {}

# ═══════════════════════════════════════════════════════════════
#  DATA LOADING
# ═══════════════════════════════════════════════════════════════
def load_excel(file):
    xl = pd.ExcelFile(file)
    dfs = {}
    for sheet in xl.sheet_names:
        try:
            df = xl.parse(sheet)
            df.columns = [str(c).strip() for c in df.columns]
            dfs[sheet] = df
        except Exception:
            pass
    return dfs

def build_data_summary(dfs):
    lines = []
    for name, df in dfs.items():
        lines.append(f"=== Sheet: '{name}' ({len(df)} rows x {len(df.columns)} cols) ===")
        lines.append(f"Columns: {list(df.columns)}")
        num_cols = df.select_dtypes(include='number').columns.tolist()
        cat_cols = df.select_dtypes(exclude='number').columns.tolist()
        lines.append(f"Numeric: {num_cols}")
        lines.append(f"Text: {cat_cols}")
        lines.append(f"Sample:\n{df.head(3).to_string()}")
        if num_cols:
            lines.append(f"Stats:\n{df[num_cols].describe().round(1).to_string()}")
        lines.append("")
    return "\n".join(lines)

# ═══════════════════════════════════════════════════════════════
#  PURE-PYTHON COLUMN FINDER
# ═══════════════════════════════════════════════════════════════
def _find_col(df, *keywords):
    for kw in keywords:
        for col in df.columns:
            if kw.lower() in col.lower().replace(" ","_").replace("-","_"):
                return col
    return None

def _mk_kpi(label, sheet, col, agg, prefix="", suffix="", desc=""):
    return {"label":label,"sheet":sheet,"column":col,"agg":agg,
            "prefix":prefix,"suffix":suffix,"found":True,"description":desc}

def _mk_chart(title, ctype, sheet, x, y, color=None):
    return {"title":title,"type":ctype,"sheet":sheet,"x":x,"y":y,"color":color,"found":True}

# ═══════════════════════════════════════════════════════════════
#  AUTO-MAPPER (pure Python, no Gemini)
# ═══════════════════════════════════════════════════════════════
def _auto_map_all_pages(dfs):
    result = {}
    sprint_sheet = sprint_df = resource_sheet = resource_df = None
    bug_sheet = bug_df = project_sheet = project_df = None

    for sheet, df in dfs.items():
        sl = sheet.lower()
        cl = " ".join(df.columns.tolist()).lower()
        if sprint_df is None and any(x in sl or x in cl for x in ["sprint","velocity","story","backlog","iteration"]):
            sprint_sheet, sprint_df = sheet, df.copy()
        if resource_df is None and any(x in sl or x in cl for x in ["resource","team","employee","headcount","utiliz","alloc"]):
            resource_sheet, resource_df = sheet, df.copy()
        if bug_df is None and any(x in sl or x in cl for x in ["bug","defect","quality","incident"]):
            bug_sheet, bug_df = sheet, df.copy()
        if project_df is None and any(x in sl or x in cl for x in ["project","portfolio","status","rag","health"]):
            project_sheet, project_df = sheet, df.copy()

    first_sheet = list(dfs.keys())[0] if dfs else None
    first_df    = dfs[first_sheet] if first_sheet else None

    # ── EXECUTIVE SUMMARY ────────────────────────────────────
    ex_kpis, ex_charts = [], []
    for sh, df in dfs.items():
        col = _find_col(df,"project_name","project","proj","initiative","programme")
        if col:
            ex_kpis.append(_mk_kpi("Total Projects",sh,col,"nunique",desc="Unique projects"))
            break
    for sh, df in dfs.items():
        col = _find_col(df,"resource_name","employee","person","member","resource","staff","name")
        if col:
            ex_kpis.append(_mk_kpi("Total Resources",sh,col,"nunique",desc="Total headcount"))
            break
    for sh, df in dfs.items():
        loc = _find_col(df,"location","shore","type","region","site","onshore","offshore")
        if loc:
            ex_kpis.append(_mk_kpi("OnShore Resources",sh,loc,"count",desc="Onshore count"))
            ex_kpis.append(_mk_kpi("OffShore Resources",sh,loc,"count",desc="Offshore count"))
            break
    df_s = sprint_df if sprint_df is not None else first_df
    sh_s = sprint_sheet if sprint_sheet is not None else first_sheet
    if df_s is not None:
        sc  = _find_col(df_s,"story_id","ticket_id","story_name","story","ticket","backlog","task_id","task","issue_id")
        sp  = _find_col(df_s,"story_point","storypoint","sp","points","estimate","effort","size")
        spc = _find_col(df_s,"sprint_name","sprint_no","sprint_number","sprint_id","sprint","iteration")
        vc  = _find_col(df_s,"velocity","completed_point","done_point","delivered","completed_sp","actual")
        if sc:  ex_kpis.append(_mk_kpi("Total Stories",sh_s,sc,"count",desc="All user stories"))
        if sp:  ex_kpis.append(_mk_kpi("Total Story Points",sh_s,sp,"sum",desc="Sum of story points"))
        if spc: ex_kpis.append(_mk_kpi("Total Sprints",sh_s,spc,"nunique",desc="Number of sprints"))
        use_v = vc or sp
        if use_v: ex_kpis.append(_mk_kpi("Avg Velocity",sh_s,use_v,"mean",desc="Avg points per sprint"))
        if spc and use_v:
            ex_charts.append(_mk_chart("Sprint Velocity Trend","line",sh_s,spc,use_v))
    if bug_df is not None:
        bid = _find_col(bug_df,"bug_id","defect_id","issue_id","id","ticket_id","number")
        if bid: ex_kpis.append(_mk_kpi("Total Bugs",bug_sheet,bid,"count",desc="All bugs"))
        env = _find_col(bug_df,"environment","env","area","level")
        if env and bid:
            ex_kpis.append(_mk_kpi("Production Bugs",bug_sheet,bid,"count",desc="Production bugs"))
            ex_kpis.append(_mk_kpi("Lower Env Bugs",bug_sheet,bid,"count",desc="Dev/QA/UAT bugs"))
    if project_df is not None:
        stat = _find_col(project_df,"status","rag","health","rag_status","colour","color","traffic","state")
        proj = _find_col(project_df,"project_name","project","proj","name")
        if stat: ex_charts.append(_mk_chart("Project Status Distribution","pie",project_sheet,stat,proj))
    result["executive_summary"] = {"kpi_mappings":ex_kpis,"chart_mappings":ex_charts,"status_mapping":{"found":False},"insights":[],"insights_request":""}

    # ── SPRINT DATA ───────────────────────────────────────────
    sp_kpis, sp_charts = [], []
    if df_s is not None:
        sc  = _find_col(df_s,"story_id","ticket_id","story_name","story","ticket","backlog","task_id","task","issue_id")
        sp  = _find_col(df_s,"story_point","storypoint","sp","points","estimate","effort","size")
        spc = _find_col(df_s,"sprint_name","sprint_no","sprint_number","sprint_id","sprint","iteration")
        vc  = _find_col(df_s,"velocity","completed_point","done_point","delivered","completed_sp","actual")
        rem = _find_col(df_s,"remaining","pending","open_point","remaining_sp","remaining_point")
        stat= _find_col(df_s,"story_status","task_status","status","state","stage")
        use_v = vc or sp
        if spc: sp_kpis.append(_mk_kpi("Total Sprints",sh_s,spc,"nunique",desc="Total sprints"))
        if sc:  sp_kpis.append(_mk_kpi("Total Stories",sh_s,sc,"count",desc="All stories"))
        if sp:  sp_kpis.append(_mk_kpi("Total Story Points",sh_s,sp,"sum",desc="Sum of story points"))
        if use_v: sp_kpis.append(_mk_kpi("Completed Points",sh_s,use_v,"sum",desc="Completed story points"))
        if rem: sp_kpis.append(_mk_kpi("Remaining Points",sh_s,rem,"sum",desc="Remaining story points"))
        if sc:  sp_kpis.append(_mk_kpi("Avg Stories Per Sprint",sh_s,sc,"count",desc="Stories per sprint"))
        if spc and sc:  sp_charts.append(_mk_chart("Stories per Sprint","bar",sh_s,spc,sc))
        if spc and sp:  sp_charts.append(_mk_chart("Story Points per Sprint","bar",sh_s,spc,sp))
        if stat:        sp_charts.append(_mk_chart("Stories by Status","pie",sh_s,stat,None))
        if spc and use_v: sp_charts.append(_mk_chart("Completion Rate by Sprint","line",sh_s,spc,use_v))
    result["sprint_data"] = {"kpi_mappings":sp_kpis,"chart_mappings":sp_charts,"status_mapping":{"found":False},"insights":[],"insights_request":""}

    # ── VELOCITY ──────────────────────────────────────────────
    vel_kpis, vel_charts = [], []
    if df_s is not None:
        spc  = _find_col(df_s,"sprint_name","sprint_no","sprint","iteration")
        vc   = _find_col(df_s,"velocity","completed_point","done_point","delivered","completed_sp","actual")
        sp   = _find_col(df_s,"story_point","storypoint","sp","points","estimate","effort")
        com  = _find_col(df_s,"committed","planned","target","planned_sp","committed_sp")
        proj = _find_col(df_s,"project_name","project","team","squad","proj")
        use_v = vc or sp
        if use_v:
            vel_kpis.append(_mk_kpi("Avg Velocity",sh_s,use_v,"mean",desc="Mean points per sprint"))
            vel_kpis.append(_mk_kpi("Peak Velocity",sh_s,use_v,"max",desc="Highest sprint velocity"))
            vel_kpis.append(_mk_kpi("Lowest Velocity",sh_s,use_v,"min",desc="Lowest sprint velocity"))
            vel_kpis.append(_mk_kpi("Total Points Delivered",sh_s,use_v,"sum",desc="Cumulative points"))
        if spc and use_v: vel_charts.append(_mk_chart("Velocity Trend","area",sh_s,spc,use_v))
        if spc and (com or use_v): vel_charts.append(_mk_chart("Committed vs Completed","bar",sh_s,spc,com or use_v))
        if proj and use_v: vel_charts.append(_mk_chart("Velocity by Project or Team","bar",sh_s,proj,use_v))
        if sp: vel_charts.append(_mk_chart("Story Point Size Distribution","histogram",sh_s,sp,None))
    result["velocity"] = {"kpi_mappings":vel_kpis,"chart_mappings":vel_charts,"status_mapping":{"found":False},"insights":[],"insights_request":""}

    # ── RESOURCES ─────────────────────────────────────────────
    res_kpis, res_charts = [], []
    df_r = resource_df if resource_df is not None else first_df
    sh_r = resource_sheet if resource_sheet is not None else first_sheet
    if df_r is not None:
        name = _find_col(df_r,"resource_name","employee","person","member","name","resource","staff")
        util = _find_col(df_r,"utiliz","alloc","capacity_used","util_percent","util_%","percentage")
        loc  = _find_col(df_r,"location","shore","onshore","offshore","type","region","site")
        proj = _find_col(df_r,"project_name","project","proj","initiative")
        team = _find_col(df_r,"team","squad","department","group")
        if name: res_kpis.append(_mk_kpi("Total Resources",sh_r,name,"nunique",desc="Total headcount"))
        if util:
            res_kpis.append(_mk_kpi("Avg Utilization %",sh_r,util,"mean",suffix="%",desc="Mean utilization"))
            res_kpis.append(_mk_kpi("Over-Allocated",sh_r,util,"count",desc="Resources >100%"))
            res_kpis.append(_mk_kpi("Bench / Available",sh_r,util,"count",desc="Resources <80%"))
        if loc: res_charts.append(_mk_chart("OnShore vs OffShore","pie",sh_r,loc,None))
        if name and util: res_charts.append(_mk_chart("Utilization by Resource or Team","bar",sh_r,name,util))
        dest = proj or team
        if dest and name: res_charts.append(_mk_chart("Resources per Project","bar",sh_r,dest,name))
        if util and dest: res_charts.append(_mk_chart("Utilization Trend","line",sh_r,dest,util))
    result["resource_utilization"] = {"kpi_mappings":res_kpis,"chart_mappings":res_charts,"status_mapping":{"found":False},"insights":[],"insights_request":""}

    # ── QUALITY ───────────────────────────────────────────────
    bug_kpis, bug_charts = [], []
    df_b = bug_df if bug_df is not None else first_df
    sh_b = bug_sheet if bug_sheet is not None else first_sheet
    if df_b is not None:
        bid  = _find_col(df_b,"bug_id","defect_id","issue_id","id","ticket_id","number")
        env  = _find_col(df_b,"environment","env","area","level","tier")
        pri  = _find_col(df_b,"priority","severity","criticality","impact")
        stat = _find_col(df_b,"status","state","resolution","stage")
        res_d= _find_col(df_b,"resolution_day","days_to_fix","days_open","fix_time","cycle_time")
        spr  = _find_col(df_b,"sprint","sprint_name","sprint_no","iteration","month","date")
        id_c = bid or (list(df_b.columns)[0] if df_b.columns.tolist() else None)
        if id_c:
            bug_kpis.append(_mk_kpi("Total Bugs",sh_b,id_c,"count",desc="All bugs"))
            bug_kpis.append(_mk_kpi("Open Bugs",sh_b,id_c,"count",desc="Open bugs"))
            bug_kpis.append(_mk_kpi("Resolved Bugs",sh_b,id_c,"count",desc="Resolved bugs"))
            bug_kpis.append(_mk_kpi("Production Bugs",sh_b,id_c,"count",desc="Production bugs"))
            bug_kpis.append(_mk_kpi("Lower Env Bugs",sh_b,id_c,"count",desc="Dev/QA/UAT bugs"))
            bug_kpis.append(_mk_kpi("Critical / P1 Bugs",sh_b,id_c,"count",desc="P1 critical bugs"))
            bug_kpis.append(_mk_kpi("Resolution Rate %",sh_b,id_c,"count",suffix="%",desc="% resolved"))
        if res_d: bug_kpis.append(_mk_kpi("Avg Resolution Days",sh_b,res_d,"mean",desc="Avg days to fix"))
        if env and id_c: bug_charts.append(_mk_chart("Bugs by Environment","bar",sh_b,env,id_c))
        if pri and id_c: bug_charts.append(_mk_chart("Bugs by Priority / Severity","pie",sh_b,pri,id_c))
        if spr and id_c: bug_charts.append(_mk_chart("Bug Trend Over Time","line",sh_b,spr,id_c))
        if stat and id_c:bug_charts.append(_mk_chart("Bugs by Status","pie",sh_b,stat,id_c))
    result["quality_metrics"] = {"kpi_mappings":bug_kpis,"chart_mappings":bug_charts,"status_mapping":{"found":False},"insights":[],"insights_request":""}

    # Status mapping
    if project_df is not None:
        stat = _find_col(project_df,"status","rag","health","rag_status","colour","color","traffic","state")
        proj = _find_col(project_df,"project_name","project","proj","name")
        if stat and proj:
            sm = {"sheet":project_sheet,"project_col":proj,"status_col":stat,"found":True}
            for pid in result:
                result[pid]["status_mapping"] = sm
    return result

# ═══════════════════════════════════════════════════════════════
#  GEMINI PAGE MAPPER
# ═══════════════════════════════════════════════════════════════
def _all_columns_index(dfs):
    lines = ["AVAILABLE SHEETS AND COLUMNS:"]
    for sheet, df in dfs.items():
        num_cols = df.select_dtypes(include='number').columns.tolist()
        cat_cols = df.select_dtypes(exclude='number').columns.tolist()
        lines.append(f'\nSheet: "{sheet}" ({len(df)} rows)')
        lines.append(f"  Numeric: {num_cols}")
        lines.append(f"  Text: {cat_cols}")
        for _, row in df.head(2).iterrows():
            lines.append(f"  Sample: {dict(row)}")
    return "\n".join(lines)

def map_page_to_data(page, data_summary, dfs):
    synonyms   = rules_cfg.get("data_rules",{}).get("synonyms",{})
    kpi_defs   = [k for s in page.get("sections",[]) if s.get("type")=="kpi_row" for k in s.get("kpis",[])]
    chart_defs = [c for s in page.get("sections",[]) if s.get("type")=="chart_row" for c in s.get("charts",[])]
    has_status  = any(s.get("type")=="status_grid"   for s in page.get("sections",[]))
    has_insight = any(s.get("type")=="insight_panel" for s in page.get("sections",[]))

    col_index    = _all_columns_index(dfs)
    synonyms_str = json.dumps(synonyms, indent=2)
    page_title   = page.get("title","")
    kpis_str     = json.dumps([{"label":k.get("label"),"hint":k.get("hint")} for k in kpi_defs], indent=2)
    charts_str   = json.dumps([{"title":c.get("title"),"type":c.get("type"),"hint":c.get("hint")} for c in chart_defs], indent=2)

    schema = '{"kpi_mappings":[{"label":"...","sheet":"...","column":"...","agg":"sum|mean|max|min|count|nunique|last","prefix":"","suffix":"","found":true}],"chart_mappings":[{"title":"...","type":"bar|line|area|pie|scatter|histogram|box","sheet":"...","x":"...","y":"...","color":null,"found":true}],"status_mapping":{"sheet":"...","project_col":"...","status_col":"...","found":true},"insights_request":"..."}'

    prompt = (
        f"{col_index}\n\nSYNONYMS:\n{synonyms_str}\n\n"
        f'PAGE: "{page_title}"\nKPIs:\n{kpis_str}\nCharts:\n{charts_str}\n'
        f"Status grid: {has_status}\nInsights: {has_insight}\n\n"
        f"Return JSON:\n{schema}\n\n"
        "Rules: use ONLY exact column names from above. Set found=false if unavailable."
    )
    raw = call_gemini(prompt, "Return ONLY valid JSON. No markdown.")
    result = parse_json(raw)
    if not isinstance(result, dict):
        return {}
    return result

# ═══════════════════════════════════════════════════════════════
#  INSIGHTS
# ═══════════════════════════════════════════════════════════════
def generate_insights(request, data_summary, page_title):
    if not request:
        return []
    schema = '[{"title":"...","body":"2-3 sentences with numbers","type":"info|warning|success|danger"}]'
    prompt = (
        f"Senior delivery analyst. Generate 3-4 insights for '{page_title}'.\n"
        f"Data:\n{data_summary[:2000]}\nFocus: {request}\n"
        f"Return JSON array:\n{schema}"
    )
    raw = call_gemini(prompt, "Return ONLY valid JSON array.")
    result = parse_json(raw)
    return result if isinstance(result, list) else []

# ═══════════════════════════════════════════════════════════════
#  KPI COMPUTATION
# ═══════════════════════════════════════════════════════════════
def fmt_val(val, prefix="", suffix=""):
    fmt = rules_cfg.get("kpi_rules",{}).get("formatting",{})
    m = fmt.get("millions_threshold",1_000_000)
    k = fmt.get("thousands_threshold",1_000)
    dp = fmt.get("decimal_places",1)
    try:
        v = float(val)
        if abs(v) >= m:    return f"{prefix}{v/m:.{dp}f}M{suffix}"
        elif abs(v) >= k:  return f"{prefix}{v/k:.{dp}f}K{suffix}"
        elif v == int(v):  return f"{prefix}{int(v):,}{suffix}"
        else:              return f"{prefix}{v:,.{dp}f}{suffix}"
    except Exception:
        return f"{prefix}{val}{suffix}"

def _do_agg(series, agg):
    if   agg=="sum":     return series.sum()
    elif agg=="mean":    return round(float(series.mean()),1)
    elif agg=="max":     return series.max()
    elif agg=="min":     return series.min()
    elif agg=="count":   return int(series.count())
    elif agg=="nunique": return int(series.nunique())
    elif agg in ("last","text"):
        d = series.dropna()
        return d.iloc[-1] if len(d) else None
    return series.sum()

def compute_kpi(mapping, dfs):
    if not mapping or not mapping.get("found",False):
        return None
    sheet  = mapping.get("sheet","")
    col    = mapping.get("column","")
    agg    = mapping.get("agg","sum")
    prefix = mapping.get("prefix","")
    suffix = mapping.get("suffix","")

    def try_compute(df, c):
        try:
            val = _do_agg(df[c], agg)
            return fmt_val(val,prefix,suffix) if val is not None else None
        except Exception:
            return None

    # Primary
    if sheet and sheet in dfs and col and col in dfs[sheet].columns:
        r = try_compute(dfs[sheet], col)
        if r is not None: return r
    # Search all sheets by exact column name
    for sh, df in dfs.items():
        if col and col in df.columns:
            r = try_compute(df, col)
            if r is not None: return r
    # Fuzzy column name match
    if col:
        col_n = col.lower().replace("_","").replace(" ","")
        for sh, df in dfs.items():
            for c in df.columns:
                if c.lower().replace("_","").replace(" ","") == col_n:
                    r = try_compute(df, c)
                    if r is not None: return r
    return None

# ═══════════════════════════════════════════════════════════════
#  CHART PALETTE & TEMPLATE
# ═══════════════════════════════════════════════════════════════
PALETTE = rules_cfg.get("chart_rules",{}).get("color_scheme",{}).get("palette",
    ["#7c3aed","#059669","#2563eb","#d97706","#dc2626","#0891b2","#ea580c","#9333ea","#16a34a","#0284c7"])
RAG = rules_cfg.get("chart_rules",{}).get("rag_colors",
    {"green":"#16a34a","amber":"#d97706","red":"#dc2626","grey":"#6b7280"})
TMPL = dict(layout=dict(
    paper_bgcolor='rgba(0,0,0,0)', plot_bgcolor='#fafafa',
    font=dict(color='#4a4a6a',family='DM Sans'),
    title=dict(font=dict(color='#1a1a2e',family='Syne',size=13)),
    xaxis=dict(gridcolor='#e8e4f0',linecolor='#d4d0e8',tickfont=dict(size=10)),
    yaxis=dict(gridcolor='#e8e4f0',linecolor='#d4d0e8',tickfont=dict(size=10)),
    legend=dict(bgcolor='rgba(255,255,255,0.8)',font=dict(color='#4a4a6a',size=10)),
    colorway=PALETTE, margin=dict(l=40,r=20,t=40,b=50),
))

# ═══════════════════════════════════════════════════════════════
#  PLOTLY CHART BUILDER
# ═══════════════════════════════════════════════════════════════
def build_chart(mapping, dfs):
    if not mapping.get("found",False): return None
    sheet = mapping.get("sheet","")
    if not sheet or sheet not in dfs: return None
    df = dfs[sheet].copy()

    x = mapping.get("x"); y = mapping.get("y"); clr = mapping.get("color")
    ct = mapping.get("type","bar"); title = mapping.get("title","")
    num = df.select_dtypes(include='number').columns.tolist()
    cat = df.select_dtypes(exclude='number').columns.tolist()

    if x and x not in df.columns: x = None
    if y and y not in df.columns: y = None
    if clr and clr not in df.columns: clr = None
    if not x: x = cat[0] if cat else (num[0] if num else None)
    if not y and num: y = next((c for c in num if c != x), num[0] if num else None)

    mb = rules_cfg.get("chart_rules",{}).get("max_bar_items",15)
    ml = rules_cfg.get("chart_rules",{}).get("max_line_points",50)
    mp = rules_cfg.get("chart_rules",{}).get("max_pie_slices",8)
    tmpl = go.layout.Template(TMPL)
    kw = dict(title=title, template=tmpl, color_discrete_sequence=PALETTE)

    def rag_colors(idx):
        return [RAG.get(str(v).lower(),PALETTE[i%len(PALETTE)]) for i,v in enumerate(idx)]

    try:
        if ct=="bar" and x and y:
            d = df.groupby(x)[y].sum().reset_index().head(mb)
            fig = px.bar(d, x=x, y=y, **kw)
            fig.update_traces(marker_line_width=0)
        elif ct=="line" and x and y:
            d = df[[x,y]].dropna().head(ml)
            fig = px.line(d, x=x, y=y, color=clr, **kw, markers=True)
        elif ct=="area" and x and y:
            d = df[[x,y]].dropna().head(ml)
            fig = px.area(d, x=x, y=y, **kw)
        elif ct=="pie":
            vc = y or (num[0] if num else None)
            if vc and x:
                d = df.groupby(x)[vc].sum().reset_index().head(mp)
                is_rag = any(str(v).lower() in ("green","amber","red") for v in d[x])
                colors = rag_colors(d[x]) if is_rag else PALETTE
                fig = px.pie(d, names=x, values=vc, title=title, template=tmpl, color_discrete_sequence=colors)
            elif x:
                d = df[x].value_counts().head(mp).reset_index()
                d.columns = [x,"count"]
                is_rag = any(str(v).lower() in ("green","amber","red") for v in d[x])
                colors = rag_colors(d[x]) if is_rag else PALETTE
                fig = px.pie(d, names=x, values="count", title=title, template=tmpl, color_discrete_sequence=colors)
            else: return None
        elif ct=="scatter" and x and y:
            fig = px.scatter(df.head(300), x=x, y=y, color=clr, **kw)
        elif ct=="histogram":
            col = x if x in num else (num[0] if num else x)
            fig = px.histogram(df, x=col, **kw)
        elif ct=="box":
            fig = px.box(df, x=x, y=y or (num[0] if num else None), **kw)
        elif ct=="heatmap" and len(num)>=2:
            corr = df[num[:8]].corr()
            fig = px.imshow(corr, text_auto=".1f", title=title, color_continuous_scale="RdYlBu", template=tmpl)
        else:
            if x and y:
                d = df.groupby(x)[y].sum().reset_index().head(mb)
                fig = px.bar(d, x=x, y=y, **kw)
            else: return None
        fig.update_layout(**TMPL["layout"])
        return fig
    except Exception:
        return None

# ═══════════════════════════════════════════════════════════════
#  MATPLOTLIB FOR PPT
# ═══════════════════════════════════════════════════════════════
def chart_to_png(mapping, dfs, w, h):
    if not mapping.get("found",False): return None
    sheet = mapping.get("sheet","")
    if not sheet or sheet not in dfs: return None
    df = dfs[sheet].copy()
    x = mapping.get("x"); y = mapping.get("y")
    ct = mapping.get("type","bar"); title = mapping.get("title","")
    num = df.select_dtypes(include='number').columns.tolist()
    cat = df.select_dtypes(exclude='number').columns.tolist()
    if x and x not in df.columns: x = None
    if y and y not in df.columns: y = None
    if not x: x = cat[0] if cat else (num[0] if num else None)
    if not y and num: y = next((c for c in num if c != x), num[0] if num else None)
    mb = rules_cfg.get("chart_rules",{}).get("max_bar_items",12)
    fig, ax = plt.subplots(figsize=(w,h), dpi=120)
    fig.patch.set_facecolor('#ffffff'); ax.set_facecolor('#fafafa')
    ax.spines[['top','right']].set_visible(False)
    ax.spines[['left','bottom']].set_color('#d4d0e8')
    ax.tick_params(colors='#4a4a6a', labelsize=8)
    ax.yaxis.grid(True, color='#e8e4f0', linewidth=0.6, zorder=0)
    ax.set_axisbelow(True)
    try:
        if ct in ("bar",) and x and y:
            d = df.groupby(x)[y].sum().head(mb)
            ax.bar(range(len(d)), d.values, color=PALETTE[:len(d)], edgecolor='white', linewidth=0.4, zorder=3)
            ax.set_xticks(range(len(d)))
            ax.set_xticklabels([str(v)[:12] for v in d.index], rotation=35, ha='right', fontsize=7)
        elif ct in ("line","area") and x and y:
            d = df[[x,y]].dropna().head(50)
            xs = range(len(d))
            ax.plot(xs, d[y], color=PALETTE[0], linewidth=2, zorder=3, marker='o', markersize=3)
            if ct=="area": ax.fill_between(xs, d[y], alpha=0.15, color=PALETTE[0])
            step = max(1,len(d)//8)
            ax.set_xticks(list(xs)[::step])
            ax.set_xticklabels([str(v)[:12] for v in d[x].iloc[::step]], rotation=35, ha='right', fontsize=7)
        elif ct=="pie" and x:
            vc = y or (num[0] if num else None)
            d = df.groupby(x)[vc].sum().head(8) if vc else df[x].value_counts().head(8)
            is_rag = any(str(v).lower() in ("green","amber","red") for v in d.index)
            colors = [RAG.get(str(v).lower(),PALETTE[i%len(PALETTE)]) for i,v in enumerate(d.index)] if is_rag else PALETTE[:len(d)]
            ax.pie(d.values, labels=d.index, colors=colors, autopct='%1.1f%%', pctdistance=0.8, textprops={'fontsize':7,'color':'#1a1a2e'})
            ax.set_facecolor('#ffffff')
        elif ct=="histogram" and (x or num):
            col = x if x in num else (num[0] if num else x)
            ax.hist(df[col].dropna(), bins=20, color=PALETTE[0], edgecolor='white', linewidth=0.4, zorder=3)
        else:
            if x and y:
                d = df.groupby(x)[y].sum().head(mb)
                ax.bar(range(len(d)), d.values, color=PALETTE[:len(d)], edgecolor='white', linewidth=0.4, zorder=3)
                ax.set_xticks(range(len(d)))
                ax.set_xticklabels([str(v)[:12] for v in d.index], rotation=35, ha='right', fontsize=7)
        ax.set_title(title, fontsize=10, fontweight='bold', color='#1a1a2e', pad=6)
        plt.tight_layout(pad=0.4)
        buf = BytesIO()
        fig.savefig(buf, format='png', dpi=120, bbox_inches='tight', facecolor='#ffffff')
        buf.seek(0); plt.close(fig)
        return buf
    except Exception:
        plt.close(fig); return None

# ═══════════════════════════════════════════════════════════════
#  STATUS GRID
# ═══════════════════════════════════════════════════════════════
def render_status_grid(mapping, dfs):
    if not mapping.get("found",False):
        st.markdown("<div class='placeholder-box'>📋 Project status data not found.</div>", unsafe_allow_html=True)
        return
    sheet = mapping.get("sheet",""); pc = mapping.get("project_col",""); sc = mapping.get("status_col","")
    if not sheet or sheet not in dfs or not pc or not sc:
        st.markdown("<div class='placeholder-box'>📋 Could not map project/status columns.</div>", unsafe_allow_html=True)
        return
    df = dfs[sheet]
    if pc not in df.columns or sc not in df.columns:
        st.markdown("<div class='placeholder-box'>📋 Columns not found in sheet.</div>", unsafe_allow_html=True)
        return
    df = df[[pc,sc]].dropna()
    cols_n = rules_cfg.get("page_rules",{}).get("section_types",{}).get("status_grid",{}).get("columns",3)
    cols = st.columns(cols_n)
    for i, (_,row) in enumerate(df.iterrows()):
        proj = str(row[pc]); status = str(row[sc]).strip().lower()
        badge = f"badge-{status}" if status in ("green","amber","red") else "badge-grey"
        emoji = {"green":"🟢","amber":"🟡","red":"🔴"}.get(status,"⚪")
        with cols[i % cols_n]:
            st.markdown(f"""<div class='status-card'>
              <div class='proj-name'>{proj}</div>
              <span class='status-badge {badge}'>{emoji} {status.capitalize()}</span>
            </div>""", unsafe_allow_html=True)

# ═══════════════════════════════════════════════════════════════
#  FUZZY MATCH
# ═══════════════════════════════════════════════════════════════
def _fuzzy_match(needle, haystack):
    if not needle or not haystack: return {}
    if needle in haystack: return haystack[needle]
    nl = needle.lower().strip()
    for k,v in haystack.items():
        if k.lower().strip() == nl: return v
    for k,v in haystack.items():
        kl = k.lower().strip()
        if nl in kl or kl in nl: return v
    nw = set(nl.split()); best_s, best_v = 0, {}
    for k,v in haystack.items():
        score = len(nw & set(k.lower().split()))
        if score > best_s: best_s,best_v = score,v
    return best_v if best_s >= 1 else {}

# ═══════════════════════════════════════════════════════════════
#  PAGE RENDERER
# ═══════════════════════════════════════════════════════════════
_chart_counter = [0]  # global counter for unique chart keys

def render_page(page, mapping, dfs):
    kpi_maps    = {m["label"]: m for m in mapping.get("kpi_mappings",[])}
    chart_maps  = {m["title"]: m for m in mapping.get("chart_mappings",[])}
    status_map  = mapping.get("status_mapping",{})
    insights    = mapping.get("insights",[])
    show_titles = rules_cfg.get("page_rules",{}).get("show_section_titles",True)
    kpi_cols    = rules_cfg.get("page_rules",{}).get("section_types",{}).get("kpi_row",{}).get("columns",4)
    chart_h     = rules_cfg.get("page_rules",{}).get("section_types",{}).get("chart_row",{}).get("height_px",380)
    show_desc   = rules_cfg.get("page_rules",{}).get("show_chart_descriptions",True)
    rendered_kpis   = set()
    rendered_charts = set()

    for section in page.get("sections",[]):
        stype = section.get("type"); stitle = section.get("title","")
        if stitle and show_titles:
            st.markdown(f"<div class='section-title'>{stitle}</div>", unsafe_allow_html=True)

        if stype == "kpi_row":
            kpi_list = section.get("kpis",[])
            render_list = []
            for kd in kpi_list:
                label = kd.get("label","")
                m = _fuzzy_match(label, kpi_maps)
                render_list.append((label, kd.get("description",""), m))
                rendered_kpis.add(label.lower())
                if m: rendered_kpis.add(m.get("label",label).lower())
            for ml,m in kpi_maps.items():
                if ml.lower() not in rendered_kpis and m.get("found"):
                    render_list.append((ml, m.get("description",""), m))
                    rendered_kpis.add(ml.lower())
            if not render_list: continue
            cols = st.columns(min(len(render_list), kpi_cols))
            for i,(label,desc,m) in enumerate(render_list):
                val = compute_kpi(m,dfs) if m else None
                with cols[i % kpi_cols]:
                    if val is not None:
                        st.markdown(f"""<div class='metric-tile'>
                          <div class='label'>{label}</div>
                          <div class='value'>{val}</div>
                          <div class='delta'>{desc[:60]}</div>
                        </div>""", unsafe_allow_html=True)
                    else:
                        st.markdown(f"""<div class='metric-tile'>
                          <div class='label'>{label}</div>
                          <div class='na'>Data not available</div>
                          <div class='delta'>{desc[:60]}</div>
                        </div>""", unsafe_allow_html=True)

        elif stype == "chart_row":
            chart_list = section.get("charts",[])
            render_list = []
            for cd in chart_list:
                ctitle = cd.get("title","")
                m = _fuzzy_match(ctitle, chart_maps)
                render_list.append((ctitle, cd.get("description",""), m))
                rendered_charts.add(ctitle.lower())
                if m: rendered_charts.add(m.get("title",ctitle).lower())
            for mt,m in chart_maps.items():
                if mt.lower() not in rendered_charts and m.get("found"):
                    render_list.append((mt,"",m))
                    rendered_charts.add(mt.lower())
            if not render_list: continue
            for row_start in range(0, len(render_list), 2):
                row_items = render_list[row_start:row_start+2]
                cols = st.columns(len(row_items))
                for col,(ctitle,cdesc,m) in zip(cols, row_items):
                    with col:
                        if m and m.get("found"):
                            fig = build_chart(m, dfs)
                            if fig:
                                if cdesc and show_desc: st.caption(cdesc)
                                _chart_counter[0] += 1
                                chart_key = f"chart_{_chart_counter[0]}_{page.get('id','p')}_{ctitle}".replace(" ","_").lower()
                                st.plotly_chart(fig, use_container_width=True,
                                                config={"displayModeBar":False},
                                                height=chart_h, key=chart_key)
                            else:
                                st.markdown(f"<div class='placeholder-box'>📊 Could not render: {ctitle}<br><small>Sheet:{m.get('sheet','')} X:{m.get('x','')} Y:{m.get('y','')}</small></div>", unsafe_allow_html=True)
                        else:
                            st.markdown(f"<div class='placeholder-box'>📊 No data mapped for:<br><b>{ctitle}</b></div>", unsafe_allow_html=True)

        elif stype == "status_grid":
            render_status_grid(status_map, dfs)

        elif stype == "insight_panel":
            icon_map = {"info":"💡","warning":"⚠️","success":"✅","danger":"🔴"}
            if insights:
                for ins in insights:
                    icon = icon_map.get(ins.get("type","info"),"💡")
                    st.markdown(f"""<div class='insight-card'>
                      <div class='ititle'>{icon} {ins.get("title","")}</div>
                      {ins.get("body","")}
                    </div>""", unsafe_allow_html=True)
            else:
                st.markdown("<div class='placeholder-box'>🤖 AI insights will appear here.</div>", unsafe_allow_html=True)

# ═══════════════════════════════════════════════════════════════
#  DASHBOARD GENERATION
# ═══════════════════════════════════════════════════════════════
def generate_dashboard(dfs):
    pages = prompt_cfg.get("pages",[])
    summary = build_data_summary(dfs)
    st.session_state.data_summary = summary
    page_data = {}
    progress = st.progress(0, text="Scanning data columns...")

    # Step 1: Pure-Python auto-mapping (always works)
    progress.progress(0.1, text="Auto-mapping columns...")
    auto_mapped = _auto_map_all_pages(dfs)

    # Step 2: Gemini enrichment
    for idx, page in enumerate(pages):
        pid = page["id"]
        frac = 0.1 + 0.7*(idx/len(pages))
        progress.progress(frac, text=f"Gemini enriching: {page.get('title','')}...")
        base = auto_mapped.get(pid, {"kpi_mappings":[],"chart_mappings":[],"status_mapping":{"found":False},"insights_request":""})
        try:
            gemini = map_page_to_data(page, summary, dfs)
            if gemini.get("kpi_mappings"):
                g_labels = {m["label"] for m in gemini["kpi_mappings"] if m.get("found")}
                extras = [m for m in base.get("kpi_mappings",[]) if m.get("label") not in g_labels]
                base["kpi_mappings"] = gemini["kpi_mappings"] + extras
            if gemini.get("chart_mappings"):
                g_titles = {m["title"] for m in gemini["chart_mappings"] if m.get("found")}
                extras = [m for m in base.get("chart_mappings",[]) if m.get("title") not in g_titles]
                base["chart_mappings"] = gemini["chart_mappings"] + extras
            if gemini.get("status_mapping",{}).get("found"):
                base["status_mapping"] = gemini["status_mapping"]
            base["insights_request"] = gemini.get("insights_request","")
        except Exception:
            pass
        page_data[pid] = base

    # Step 3: AI insights
    progress.progress(0.85, text="Generating AI insights...")
    for page in pages:
        pid = page["id"]
        if any(s.get("type")=="insight_panel" for s in page.get("sections",[])):
            try:
                page_data[pid]["insights"] = generate_insights(
                    page_data[pid].get("insights_request",""), summary, page.get("title",""))
            except Exception:
                page_data[pid]["insights"] = []

    progress.progress(1.0, text="Dashboard ready!")
    progress.empty()
    st.session_state.page_data = page_data
    st.session_state.dashboard_ready = True

# ═══════════════════════════════════════════════════════════════
#  CHAT
# ═══════════════════════════════════════════════════════════════
def _detect_intent(msg):
    intents = ["add_kpi","remove_kpi","add_chart","remove_chart","change_chart","question"]
    pages_str = str([p.get("title") for p in prompt_cfg.get("pages",[])])
    prompt = (f"Classify this message into one intent.\nDashboard pages:{pages_str}\n"
              f"Message:\"{msg}\"\nIntents:{intents}\nReply with ONLY the intent key.")
    result = call_gemini(prompt,"Reply with only the intent key.").strip().lower()
    for key in intents:
        if key in result: return key
    return "question"

def _resolve_page_id(msg):
    pages = prompt_cfg.get("pages",[])
    ml = msg.lower()
    for p in pages:
        if p.get("title","").lower() in ml: return p["id"]
    for p in pages:
        if any(w in ml for w in p.get("title","").lower().split() if len(w)>3): return p["id"]
    return pages[0]["id"] if pages else None

def _ensure_live_pages():
    if not st.session_state.live_pages:
        st.session_state.live_pages = {p["id"]: copy.deepcopy(p) for p in prompt_cfg.get("pages",[])}

def _apply_add_kpi(msg, page_id, dfs):
    col_index = _all_columns_index(dfs)
    schema = '{"label":"...","sheet":"...","column":"...","agg":"sum|mean|max|min|count|nunique|last","prefix":"","suffix":"","found":true,"description":"..."}'
    prompt = f"{col_index}\n\nUser wants to add KPI: \"{msg}\"\nReturn JSON:\n{schema}\nUse only exact column names."
    raw = call_gemini(prompt,"Return ONLY valid JSON.")
    kpi_map = parse_json(raw)
    if not kpi_map or not kpi_map.get("found"):
        return "Could not find matching data for that KPI. Try being more specific about the column name."
    _ensure_live_pages()
    pid = page_id
    page_data = st.session_state.page_data
    page_data.setdefault(pid,{"kpi_mappings":[],"chart_mappings":[],"status_mapping":{"found":False},"insights":[]})
    page_data[pid].setdefault("kpi_mappings",[]).append(kpi_map)
    lp = st.session_state.live_pages
    if pid in lp:
        page = lp[pid]
        ks = [s for s in page.get("sections",[]) if s.get("type")=="kpi_row"]
        if ks: ks[-1].setdefault("kpis",[]).append({"label":kpi_map.get("label","New KPI"),"description":kpi_map.get("description",""),"hint":""})
        else: page.setdefault("sections",[]).append({"type":"kpi_row","title":"Added Metrics","kpis":[{"label":kpi_map.get("label","New KPI"),"description":kpi_map.get("description",""),"hint":""}]})
    st.session_state.page_data = page_data
    val = compute_kpi(kpi_map,dfs)
    return f"✅ Added **{kpi_map.get('label','KPI')}** tile — value: **{val or 'N/A'}**"

def _apply_remove_kpi(msg, page_id):
    _ensure_live_pages()
    lp = st.session_state.live_pages
    page = lp.get(page_id,{})
    all_labels = [k.get("label","") for s in page.get("sections",[]) if s.get("type")=="kpi_row" for k in s.get("kpis",[])]
    if not all_labels: return "No KPI tiles found on that page."
    prompt = f"User said:\"{msg}\"\nKPI labels:{all_labels}\nWhich label to remove? Reply with exact label only."
    target = call_gemini(prompt,"Reply with exact label only.").strip().strip('"').strip("'")
    removed = False
    for s in page.get("sections",[]):
        if s.get("type")=="kpi_row":
            before = len(s.get("kpis",[]))
            s["kpis"] = [k for k in s.get("kpis",[]) if k.get("label","").lower()!=target.lower()]
            if len(s.get("kpis",[]))<before: removed=True
    pd_ = st.session_state.page_data
    if page_id in pd_:
        pd_[page_id]["kpi_mappings"] = [m for m in pd_[page_id].get("kpi_mappings",[]) if m.get("label","").lower()!=target.lower()]
    st.session_state.live_pages = lp; st.session_state.page_data = pd_
    return f"✅ Removed **{target}**." if removed else f"Couldn't find '{target}'. Available:{all_labels}"

def _apply_add_chart(msg, page_id, dfs):
    col_index = _all_columns_index(dfs)
    schema = '{"title":"...","type":"bar|line|area|pie|scatter|histogram","sheet":"...","x":"...","y":"...","color":null,"found":true,"description":"..."}'
    prompt = f"{col_index}\n\nUser wants chart: \"{msg}\"\nReturn JSON:\n{schema}"
    raw = call_gemini(prompt,"Return ONLY valid JSON.")
    cm = parse_json(raw)
    if not cm or not cm.get("found"): return "Couldn't find matching data for that chart."
    _ensure_live_pages()
    lp = st.session_state.live_pages; page = lp.get(page_id,{})
    cs = [s for s in page.get("sections",[]) if s.get("type")=="chart_row"]
    cdef = {"title":cm.get("title","New Chart"),"type":cm.get("type","bar"),"description":cm.get("description",""),"hint":""}
    if cs and len(cs[-1].get("charts",[]))<2: cs[-1].setdefault("charts",[]).append(cdef)
    else: page.setdefault("sections",[]).append({"type":"chart_row","title":"Added Charts","charts":[cdef]})
    pd_ = st.session_state.page_data
    pd_.setdefault(page_id,{"kpi_mappings":[],"chart_mappings":[],"status_mapping":{"found":False},"insights":[]})
    pd_[page_id].setdefault("chart_mappings",[]).append(cm)
    st.session_state.live_pages=lp; st.session_state.page_data=pd_
    return f"✅ Added **{cm.get('title','New Chart')}** ({cm.get('type','bar')} chart)."

def _apply_remove_chart(msg, page_id):
    _ensure_live_pages()
    lp = st.session_state.live_pages; page = lp.get(page_id,{})
    all_titles = [c.get("title","") for s in page.get("sections",[]) if s.get("type")=="chart_row" for c in s.get("charts",[])]
    if not all_titles: return "No charts found on that page."
    prompt = f"User said:\"{msg}\"\nChart titles:{all_titles}\nWhich to remove? Reply with exact title only."
    target = call_gemini(prompt,"Reply with exact title only.").strip().strip('"').strip("'")
    removed=False
    for s in page.get("sections",[]):
        if s.get("type")=="chart_row":
            before=len(s.get("charts",[])); s["charts"]=[c for c in s.get("charts",[]) if c.get("title","").lower()!=target.lower()]
            if len(s.get("charts",[]))<before: removed=True
    pd_=st.session_state.page_data
    if page_id in pd_: pd_[page_id]["chart_mappings"]=[m for m in pd_[page_id].get("chart_mappings",[]) if m.get("title","").lower()!=target.lower()]
    st.session_state.live_pages=lp; st.session_state.page_data=pd_
    return f"✅ Removed **{target}**." if removed else f"Couldn't find '{target}'. Available:{all_titles}"

def _apply_change_chart(msg, page_id):
    _ensure_live_pages()
    lp = st.session_state.live_pages; page = lp.get(page_id,{})
    all_titles = [c.get("title","") for s in page.get("sections",[]) if s.get("type")=="chart_row" for c in s.get("charts",[])]
    if not all_titles: return "No charts found on that page."
    chart_types=["bar","line","area","pie","scatter","histogram","box","heatmap"]
    schema='{"target_title":"exact title","new_type":"bar|line|area|pie|scatter|histogram|box|heatmap"}'
    prompt=f"User said:\"{msg}\"\nCharts:{all_titles}\nTypes:{chart_types}\nReturn JSON:{schema}"
    raw=call_gemini(prompt,"Return ONLY valid JSON.")
    result=parse_json(raw)
    if not result: return "Couldn't understand. Try: 'Change Velocity Trend to bar chart'."
    target=result.get("target_title",""); new_type=result.get("new_type","bar")
    changed_page=False
    for s in page.get("sections",[]):
        if s.get("type")=="chart_row":
            for c in s.get("charts",[]):
                if c.get("title","").lower()==target.lower(): c["type"]=new_type; changed_page=True
    pd_=st.session_state.page_data; changed_map=False
    if page_id in pd_:
        for m in pd_[page_id].get("chart_mappings",[]):
            if m.get("title","").lower()==target.lower(): m["type"]=new_type; changed_map=True
    st.session_state.live_pages=lp; st.session_state.page_data=pd_
    return f"✅ Changed **{target}** to **{new_type}** chart." if (changed_page or changed_map) else f"Couldn't find '{target}'."

def _answer_question(msg):
    history="\n".join(f"{'User' if m['role']=='user' else 'AI'}: {m['content']}" for m in st.session_state.chat_history[-4:])
    summary=st.session_state.get("data_summary","")[:2000]
    prompt=(f"Dashboard:{dash_meta.get('title')}\nData:\n{summary}\nChat:\n{history}\nQuestion:{msg}")
    return call_gemini(prompt,"You are a helpful dashboard analyst. Answer in 2-4 sentences.")

def handle_chat(msg, dfs):
    _ensure_live_pages()
    intent  = _detect_intent(msg)
    page_id = _resolve_page_id(msg)
    if intent=="add_kpi":      return _apply_add_kpi(msg,page_id,dfs)
    elif intent=="remove_kpi": return _apply_remove_kpi(msg,page_id)
    elif intent=="add_chart":  return _apply_add_chart(msg,page_id,dfs)
    elif intent=="remove_chart":return _apply_remove_chart(msg,page_id)
    elif intent=="change_chart":return _apply_change_chart(msg,page_id)
    else:                      return _answer_question(msg)

# ═══════════════════════════════════════════════════════════════
#  PPT EXPORT
# ═══════════════════════════════════════════════════════════════
def export_pptx(dfs):
    exp=rules_cfg.get("export_rules",{}).get("pptx",{}); company=exp.get("company_name","")
    PURPLE=RGBColor(0x7C,0x3A,0xED); LP=RGBColor(0xA7,0x8B,0xFA)
    DARK=RGBColor(0x0A,0x0A,0x0F); CARD=RGBColor(0x13,0x13,0x1F)
    WHITE=RGBColor(0xFF,0xFF,0xFF); GRAY=RGBColor(0x6B,0x6B,0x8A)
    prs=Presentation(); prs.slide_width=Inches(13.33); prs.slide_height=Inches(7.5)
    def blank():
        sl=prs.slides.add_slide(prs.slide_layouts[6])
        bg=sl.background.fill; bg.solid(); bg.fore_color.rgb=DARK
        return sl
    def tb(sl,text,l,t,w,h,sz=14,bold=False,color=WHITE,align=PP_ALIGN.LEFT,italic=False):
        box=sl.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
        tf=box.text_frame; tf.word_wrap=True; p=tf.paragraphs[0]; p.alignment=align
        run=p.add_run(); run.text=str(text); run.font.name='Calibri'; run.font.size=Pt(sz)
        run.font.bold=bold; run.font.italic=italic; run.font.color.rgb=color
    def rect(sl,l,t,w,h,fill,line=None):
        sh=sl.shapes.add_shape(1,Inches(l),Inches(t),Inches(w),Inches(h))
        sh.fill.solid(); sh.fill.fore_color.rgb=fill
        if line: sh.line.color.rgb=line
        else: sh.line.fill.background()
    sl=blank(); rect(sl,0,2.8,0.08,2.0,PURPLE)
    tb(sl,dash_meta.get("title","Dashboard"),0.4,2.9,10,0.9,42,True)
    tb(sl,dash_meta.get("subtitle",""),0.4,4.0,9,0.5,17,color=GRAY)
    footer=f"{company}  ·  Powered by Gemini AI" if company else "Powered by Gemini AI"
    tb(sl,footer,0.4,6.9,10,0.35,9,color=GRAY,italic=True)
    lp=st.session_state.get("live_pages",{p["id"]:p for p in prompt_cfg.get("pages",[])})
    pd_=st.session_state.page_data
    pages_to_export=exp.get("include_pages",[p["id"] for p in prompt_cfg.get("pages",[])])
    for pid in pages_to_export:
        if pid not in lp or pid not in pd_: continue
        page=lp[pid]; mapping=pd_[pid]
        chart_maps={m["title"]:m for m in mapping.get("chart_mappings",[])}
        kpi_maps={m["label"]:m for m in mapping.get("kpi_mappings",[])}
        sl=blank()
        tb(sl,f"{page.get('icon','')} {page.get('title','')}",0.4,0.2,12,0.45,10,True,LP)
        tb(sl,dash_meta.get("title",""),0.4,0.6,12,0.45,20,True)
        kpi_defs=[k for s in page.get("sections",[]) if s.get("type")=="kpi_row" for k in s.get("kpis",[])]
        cw,ch,xs,ys,gap=2.8,1.35,0.3,1.2,0.22
        for i,kd in enumerate(kpi_defs[:8]):
            r,c_=i//4,i%4; x=xs+c_*(cw+gap); y=ys+r*(ch+0.25)
            rect(sl,x,y,cw,ch,CARD,RGBColor(0x2A,0x2A,0x3E)); rect(sl,x,y,0.05,ch,PURPLE)
            tb(sl,kd.get("label",""),x+0.12,y+0.1,cw-0.2,0.28,8,True,GRAY)
            m=_fuzzy_match(kd.get("label",""),kpi_maps)
            val=compute_kpi(m,dfs) if m else None
            tb(sl,val or "N/A",x+0.12,y+0.42,cw-0.2,0.55,22,True,WHITE if val else GRAY)
            tb(sl,kd.get("description","")[:50],x+0.12,y+1.02,cw-0.2,0.28,7,color=GRAY)
        chart_defs=[c for s in page.get("sections",[]) if s.get("type")=="chart_row" for c in s.get("charts",[])]
        for i in range(0,len(chart_defs),2):
            sl=blank(); tb(sl,f"{page.get('title','')} — Charts",0.3,0.15,12,0.35,9,True,LP)
            batch=chart_defs[i:i+2]
            pos=[(0.3,0.6,6.1,6.4),(6.7,0.6,6.1,6.4)] if len(batch)==2 else [(0.3,0.6,12.5,6.4)]
            for j,cd in enumerate(batch):
                lf,tp_,wd,ht=pos[j]
                cm=_fuzzy_match(cd.get("title",""),chart_maps)
                buf=chart_to_png(cm,dfs,wd*0.88,ht*0.76) if cm else None
                if buf: sl.shapes.add_picture(buf,Inches(lf),Inches(tp_+0.42),Inches(wd),Inches(ht-0.42))
                tb(sl,cd.get("title",""),lf,tp_,wd,0.38,11,True)
        ins_list=mapping.get("insights",[])
        if ins_list:
            sl=blank(); tb(sl,f"{page.get('title','')} — Insights",0.3,0.15,12,0.35,9,True,LP)
            tb(sl,"AI-Generated Observations",0.3,0.5,12,0.45,20,True)
            for i,ins in enumerate(ins_list[:4]):
                r,c_=i//2,i%2; x=0.3+c_*6.5; y=1.2+r*2.7
                rect(sl,x,y,6.1,2.3,CARD,RGBColor(0x2A,0x2A,0x3E)); rect(sl,x,y,6.1,0.05,PURPLE)
                tb(sl,ins.get("title",""),x+0.18,y+0.14,5.7,0.4,11,True)
                tb(sl,ins.get("body",""),x+0.18,y+0.58,5.7,1.5,9,color=GRAY)
    buf=BytesIO(); prs.save(buf); buf.seek(0)
    return buf

# ═══════════════════════════════════════════════════════════════
#  SIDEBAR
# ═══════════════════════════════════════════════════════════════
with st.sidebar:
    st.markdown(f"### {dash_meta.get('logo_emoji','📊')} {dash_meta.get('title','Dashboard')}")
    st.markdown(f"<div style='color:#8b8ba8;font-size:.76rem;margin-bottom:.5rem;'>{dash_meta.get('subtitle','')}</div>", unsafe_allow_html=True)
    st.markdown("---")
    st.markdown("#### 🔑 Gemini API Key")
    api_key = st.text_input("", type="password", placeholder="AIza…", value=st.session_state.api_key, label_visibility="collapsed")
    if api_key:
        st.session_state.api_key = api_key
        st.success("API key set", icon="✅")
    st.markdown("---")
    st.markdown("#### 📂 Upload Excel Data")
    uploaded = st.file_uploader("", type=["xlsx","xls"], label_visibility="collapsed")
    if uploaded and not st.session_state.file_uploaded:
        with st.spinner("Reading sheets…"):
            dfs = load_excel(uploaded)
            st.session_state.dataframes = dfs
            st.session_state.sheet_names = list(dfs.keys())
            st.session_state.file_uploaded = True
            st.session_state.dashboard_ready = False
            st.session_state.page_data = {}
            st.session_state.live_pages = {}
        st.success(f"Loaded {len(dfs)} sheet(s)", icon="📊")
    if st.session_state.file_uploaded:
        if st.button("🔄 Reset / New File"):
            for k in ["file_uploaded","dashboard_ready","page_data","dataframes","sheet_names","chat_history","live_pages"]:
                st.session_state[k] = {} if k in ("page_data","dataframes","live_pages") else ([] if k in ("sheet_names","chat_history") else False)
            st.rerun()
        st.markdown("**Sheets:**")
        for s in st.session_state.sheet_names:
            st.markdown(f"  `{s}` — {len(st.session_state.dataframes[s])} rows")
    st.markdown("---")
    if st.session_state.file_uploaded and st.session_state.api_key:
        if st.button("🚀 Generate Dashboard", use_container_width=True):
            _chart_counter[0] = 0
            with st.spinner("Gemini is mapping your data…"):
                generate_dashboard(st.session_state.dataframes)
            st.rerun()
    if st.session_state.dashboard_ready:
        st.markdown("---")
        if st.button("📥 Export to PowerPoint", use_container_width=True):
            with st.spinner("Building presentation…"):
                pptx_buf = export_pptx(st.session_state.dataframes)
            st.download_button("⬇️ Download .pptx", data=pptx_buf, file_name="dashboard.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation", use_container_width=True)
        if st.button("↺ Reset Chat Changes", use_container_width=True):
            st.session_state.live_pages = {p["id"]:copy.deepcopy(p) for p in prompt_cfg.get("pages",[])}
            st.session_state.chat_history = []
            st.rerun()
    st.markdown("---")
    st.markdown("""<div style='font-size:.73rem;color:#8b8ba8;line-height:1.8;'>
    <b>Config files:</b><br>📄 dashboard_prompt.yaml<br>📋 rules.yaml</div>""", unsafe_allow_html=True)

# ═══════════════════════════════════════════════════════════════
#  MAIN
# ═══════════════════════════════════════════════════════════════
if not st.session_state.file_uploaded:
    st.markdown(f"""<div style='text-align:center;padding:3.5rem 1rem 2rem;'>
      <div style='font-family:Syne,sans-serif;font-size:2.6rem;font-weight:800;
                  background:linear-gradient(135deg,#6d28d9,#059669);
                  -webkit-background-clip:text;-webkit-text-fill-color:transparent;'>
        {dash_meta.get("title","AI Dashboard Studio")}</div>
      <div style='color:#6b7280;font-size:1rem;margin-bottom:2rem;'>
        {dash_meta.get("subtitle","Config-driven · Gemini-powered · Excel-ready")}</div>
    </div>""", unsafe_allow_html=True)
    c1,c2,c3 = st.columns(3)
    for col,icon,title,body in [
        (c1,"📄","Config-Driven","Define pages and KPIs in <b>dashboard_prompt.yaml</b>. No code needed."),
        (c2,"🤖","Smart Mapping","Gemini maps your Excel columns to every KPI and chart automatically."),
        (c3,"📤","Export to PPT","One-click PowerPoint export with all charts and AI insights."),
    ]:
        with col:
            st.markdown(f"<div class='insight-card'><div class='ititle'>{icon} {title}</div>{body}</div>", unsafe_allow_html=True)
    st.markdown("---")
    st.markdown("#### Pages defined in `dashboard_prompt.yaml`")
    pages = prompt_cfg.get("pages",[])
    if pages:
        pcols = st.columns(min(len(pages),5))
        for i,page in enumerate(pages):
            nk=sum(len(s.get("kpis",[])) for s in page.get("sections",[]) if s.get("type")=="kpi_row")
            nc=sum(len(s.get("charts",[])) for s in page.get("sections",[]) if s.get("type")=="chart_row")
            with pcols[i%5]:
                st.markdown(f"<div class='insight-card'><div class='ititle'>{page.get('icon','')} {page.get('title','')}</div><span style='color:#8b8ba8;font-size:.76rem;'>{nk} KPIs · {nc} Charts</span></div>", unsafe_allow_html=True)

elif not st.session_state.dashboard_ready:
    st.markdown("""<div style='text-align:center;padding:2.5rem;'>
      <div style='font-size:1.3rem;font-family:Syne,sans-serif;color:#6d28d9;margin-bottom:.7rem;'>Data loaded — ready to generate</div>
      <div style='color:#6b7280;'>Click <b>Generate Dashboard</b> in the sidebar</div>
    </div>""", unsafe_allow_html=True)
    st.markdown("#### Data Preview")
    tabs = st.tabs(st.session_state.sheet_names)
    for tab,name in zip(tabs, st.session_state.sheet_names):
        with tab:
            st.dataframe(st.session_state.dataframes[name].head(15), use_container_width=True, height=280)

else:
    # Reset chart counter on each render
    _chart_counter[0] = 0

    hc1,hc2 = st.columns([9,2])
    with hc1:
        st.markdown(f"""<div style='font-family:Syne,sans-serif;font-size:1.7rem;font-weight:800;color:#1a1a2e;margin-bottom:.1rem;'>
          {dash_meta.get("title","Dashboard")}</div>
          <div style='color:#6b7280;font-size:.86rem;'>{dash_meta.get("subtitle","")}</div>""", unsafe_allow_html=True)
    with hc2:
        st.markdown("""<div style='text-align:right;padding-top:.3rem;'>
          <span style='background:#f0ebff;border:1px solid #c4b5fd;border-radius:20px;
                       padding:.22rem .65rem;font-size:.7rem;color:#6d28d9;
                       font-family:Syne,sans-serif;font-weight:700;'>✦ Gemini Powered</span>
        </div>""", unsafe_allow_html=True)
    st.markdown("---")

    _ensure_live_pages()
    pages = prompt_cfg.get("pages",[])
    live_pages = st.session_state.live_pages
    tabs = st.tabs([f"{live_pages.get(p['id'],p).get('icon','')} {live_pages.get(p['id'],p).get('title','')}" for p in pages])
    for tab,page in zip(tabs,pages):
        pid = page["id"]
        with tab:
            render_page(live_pages.get(pid,page), st.session_state.page_data.get(pid,{}), st.session_state.dataframes)

    st.markdown("---")
    st.markdown("<div class='section-title'>💬 Modify Dashboard via Chat</div>", unsafe_allow_html=True)
    st.markdown("""<div style='background:#f0ebff;border:1px solid #e0d9ff;border-radius:10px;
        padding:.7rem 1rem;margin-bottom:.8rem;font-size:.8rem;color:#4a3080;line-height:1.8;'>
        <b>What you can say:</b><br>
        ➕ <i>"Add a KPI for total resolved bugs"</i> &nbsp;·&nbsp;
        ➕ <i>"Add a bar chart of bugs by priority"</i><br>
        🔄 <i>"Change Sprint Velocity Trend to a bar chart"</i><br>
        ❌ <i>"Remove the Avg Velocity tile"</i> &nbsp;·&nbsp;
        ❌ <i>"Remove the Utilization Trend chart"</i><br>
        💬 <i>"Which project has the most bugs?"</i>
    </div>""", unsafe_allow_html=True)

    for msg_item in st.session_state.chat_history[-10:]:
        css = "chat-user" if msg_item["role"]=="user" else "chat-ai"
        icon = "👤" if msg_item["role"]=="user" else "✦"
        st.markdown(f"<div class='{css}'><b>{icon}</b> {msg_item['content']}</div>", unsafe_allow_html=True)

    user_input = st.chat_input("Add a tile, change a chart, remove something, or ask a question…")
    if user_input:
        st.session_state.chat_history.append({"role":"user","content":user_input})
        with st.spinner("Applying…"):
            reply = handle_chat(user_input, st.session_state.dataframes)
        st.session_state.chat_history.append({"role":"assistant","content":reply})
        st.rerun()

    with st.expander("🗃️ Raw Data Explorer"):
        tabs2 = st.tabs(st.session_state.sheet_names)
        for tab2,name in zip(tabs2, st.session_state.sheet_names):
            with tab2:
                st.dataframe(st.session_state.dataframes[name], use_container_width=True, height=240)

    with st.expander("🔍 Mapping Debugger"):
        st.markdown("<div style='color:#6b7280;font-size:.8rem;'>Use this to diagnose blank tiles. If Found=❌ the column name doesn't match your data.</div>", unsafe_allow_html=True)
        st.markdown("**Your Excel columns:**")
        for sheet,df in st.session_state.dataframes.items():
            num_cols=df.select_dtypes(include='number').columns.tolist()
            cat_cols=df.select_dtypes(exclude='number').columns.tolist()
            cc1,cc2=st.columns(2)
            with cc1: st.markdown(f"**`{sheet}`** Numeric: `{num_cols}`")
            with cc2: st.markdown(f"Text: `{cat_cols}`")
        st.markdown("---")
        for page in prompt_cfg.get("pages",[]):
            pid=page["id"]; mapping=st.session_state.page_data.get(pid,{})
            with st.expander(f"{page.get('icon','')} {page.get('title','')}"):
                kpi_maps_dbg=mapping.get("kpi_mappings",[])
                chart_maps_dbg=mapping.get("chart_mappings",[])
                if kpi_maps_dbg:
                    st.markdown("**KPI Mappings:**")
                    rows=[{"Label":m.get("label",""),"Sheet":m.get("sheet","") or "❌","Column":m.get("column","") or "❌","Agg":m.get("agg",""),"Found":"✅" if m.get("found") else "❌","Value":compute_kpi(m,st.session_state.dataframes) or "N/A"} for m in kpi_maps_dbg]
                    st.dataframe(pd.DataFrame(rows), use_container_width=True, hide_index=True)
                if chart_maps_dbg:
                    st.markdown("**Chart Mappings:**")
                    rows=[{"Title":m.get("title",""),"Type":m.get("type",""),"Sheet":m.get("sheet","") or "❌","X":m.get("x","") or "❌","Y":m.get("y","") or "❌","Found":"✅" if m.get("found") else "❌"} for m in chart_maps_dbg]
                    st.dataframe(pd.DataFrame(rows), use_container_width=True, hide_index=True)